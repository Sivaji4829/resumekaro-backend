from fastapi import FastAPI, HTTPException, Depends, Header, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import os
import httpx
import asyncio
import math
import io
import json
from supabase import create_client, Client
from dotenv import load_dotenv

# Load variables from the .env file
load_dotenv()

# --- Configuration & Setup ---
# Securely load from environment variables. 
# NO HARDCODED SECRETS HERE FOR GITHUB SECURITY.
SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY", "")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")

app = FastAPI(title="ResumeKaro API")

# Allow frontend to communicate with backend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Pydantic Models ---
class ResumeContent(BaseModel):
    personal: Dict[str, Any]
    summary: Optional[str] = ""
    experience: Optional[List[Dict[str, Any]]] = []
    education: Optional[List[Dict[str, Any]]] = []
    skills: Optional[str] = ""
    links: Optional[List[Dict[str, Any]]] = []
    projects: Optional[List[Dict[str, Any]]] = []
    certifications: Optional[List[Dict[str, Any]]] = []
    achievements: Optional[str] = ""
    languages: Optional[str] = ""
    is_deleted: Optional[bool] = False 
    ats_used: Optional[bool] = False   

class AtsAnalyzeRequest(BaseModel):
    resume_data: ResumeContent
    job_description: str


# --- Security & Limits Dependencies ---
def get_auth_context(authorization: str = Header(None)):
    """Verifies the JWT token and creates a Supabase client injected with the user's session to bypass RLS limits"""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid Authorization header")
    
    token = authorization.split(" ")[1]
    try:
        base_client = create_client(SUPABASE_URL, SUPABASE_KEY)
        user_response = base_client.auth.get_user(token)
        
        if not user_response.user:
            raise HTTPException(status_code=401, detail="Invalid token")
            
        auth_client = create_client(SUPABASE_URL, SUPABASE_KEY)
        auth_client.postgrest.auth(token)
        
        return {
            "client": auth_client,
            "user_id": user_response.user.id
        }
    except Exception as e:
        raise HTTPException(status_code=401, detail=f"Authentication failed: {str(e)}")

def check_usage_limits(auth_context: dict = Depends(get_auth_context)):
    """Strictly enforces free tier limits for ATS Scans: exactly 1 lifetime scan allowed."""
    client = auth_context["client"]
    user_id = auth_context["user_id"]
    
    try:
        resumes = client.table('resumes').select('id, content').eq('user_id', user_id).execute()
        
        if len(resumes.data) == 0:
            raise HTTPException(status_code=400, detail="Please save your profile to the database first before running an ATS scan.")
        
        for r in resumes.data:
            content = r.get('content', {})
            if content.get('ats_used', False) == True:
                raise HTTPException(
                    status_code=403, 
                    detail="LIFETIME LIMIT REACHED: You have already used your 1 free ATS scan. Please upgrade to Premium for unlimited scans."
                )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error checking limits: {e}")
        
    return auth_context


# --- Utility Functions ---
def cosine_similarity(v1, v2):
    dot_product = sum(x * y for x, y in zip(v1, v2))
    magnitude1 = math.sqrt(sum(x * x for x in v1))
    magnitude2 = math.sqrt(sum(x * x for x in v2))
    if magnitude1 == 0 or magnitude2 == 0:
        return 0
    return dot_product / (magnitude1 * magnitude2)


# --- Routes ---

@app.post("/api/v1/resumes")
async def save_resume(resume_payload: dict, auth_context: dict = Depends(get_auth_context)):
    """Strictly controlled endpoint to save/update a resume"""
    client = auth_context["client"]
    user_id = auth_context["user_id"]
    resume_id = resume_payload.get("id")
    
    resumes = client.table('resumes').select('id, content').eq('user_id', user_id).execute()
    
    if not resume_id:
        if len(resumes.data) >= 1:
            raise HTTPException(
                status_code=403, 
                detail="LIFETIME LIMIT REACHED: Free tier is restricted to 1 resume creation. Upgrade to Premium to create multiple resumes."
            )
    else:
        user_resume_ids = [r['id'] for r in resumes.data]
        if resume_id not in user_resume_ids:
            raise HTTPException(status_code=403, detail="Unauthorized to modify this resume.")
            
        for r in resumes.data:
            if r['id'] == resume_id:
                if r.get('content', {}).get('ats_used', False):
                    resume_payload['content']['ats_used'] = True
    
    resume_payload["user_id"] = user_id
    
    try:
        res = client.table('resumes').upsert(resume_payload).execute()
        return res.data
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/v1/resumes/{resume_id}")
async def delete_resume(resume_id: str, auth_context: dict = Depends(get_auth_context)):
    """Performs a SOFT DELETE. Hides the resume from UI but keeps the DB record to enforce the 1-resume lifetime limit."""
    client = auth_context["client"]
    user_id = auth_context["user_id"]
    
    try:
        resume = client.table('resumes').select('content').eq('id', resume_id).eq('user_id', user_id).execute()
        if not resume.data:
            raise HTTPException(status_code=404, detail="Resume not found")
            
        content = resume.data[0].get('content', {})
        content['is_deleted'] = True
        
        client.table('resumes').update({"content": content}).eq('id', resume_id).execute()
        return {"message": "Resume successfully deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/v1/limits/check")
async def check_account_limits(auth_context: dict = Depends(get_auth_context)):
    """Endpoint for the frontend to proactively check if the user has hit their free tier lifetime limits"""
    client = auth_context["client"]
    user_id = auth_context["user_id"]
    
    try:
        resumes = client.table('resumes').select('id, content').eq('user_id', user_id).execute()
        resume_count = len(resumes.data)
        
        scan_limit_reached = False
        for r in resumes.data:
            content = r.get('content', {})
            if content.get('ats_used', False) == True:
                scan_limit_reached = True
                break
            
        return {
            "resume_limit_reached": resume_count >= 1, 
            "scan_limit_reached": scan_limit_reached,  
            "current_resumes": resume_count,
            "premium_active": False 
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to check limits: {str(e)}")


@app.post("/api/v1/utils/extract-text")
async def extract_text_from_file(file: UploadFile = File(...)):
    """Extracts text from uploaded PDF, DOCX, PPTX, or TXT files"""
    content = await file.read()
    text = ""
    filename = file.filename.lower()
    
    try:
        if filename.endswith(".pdf"):
            import pypdf
            pdf = pypdf.PdfReader(io.BytesIO(content))
            text = " ".join([page.extract_text() for page in pdf.pages if page.extract_text()])
        elif filename.endswith(".docx") or filename.endswith(".doc"):
            import docx
            doc = docx.Document(io.BytesIO(content))
            text = " ".join([paragraph.text for paragraph in doc.paragraphs])
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            import pptx
            ppt = pptx.Presentation(io.BytesIO(content))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
        else:
            text = content.decode("utf-8")
            
        if not text.strip():
            raise ValueError("No readable text found in document.")
            
        return {"text": text.strip()}
    except ImportError as e:
        raise HTTPException(status_code=500, detail=f"Server missing required library. Please run: pip install pypdf python-docx python-pptx. Error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse file: {str(e)}")


@app.post("/api/v1/ats/score-fast")
async def score_fast_embeddings(request: AtsAnalyzeRequest, auth_context: dict = Depends(check_usage_limits)):
    client = auth_context["client"]
    user_id = auth_context["user_id"]
    url = f"https://generativelanguage.googleapis.com/v1beta/models/text-embedding-004:batchEmbedContents?key={GEMINI_API_KEY}"
    
    resume_text = str(request.resume_data.model_dump())[:8000]
    jd_text = request.job_description[:8000]
    payload = {
        "requests": [
            {"model": "models/text-embedding-004", "content": {"parts": [{"text": resume_text}]}},
            {"model": "models/text-embedding-004", "content": {"parts": [{"text": jd_text}]}}
        ]
    }

    try:
        async with httpx.AsyncClient() as client_http:
            response = await client_http.post(url, headers={"Content-Type": "application/json"}, json=payload, timeout=60.0)
            if response.status_code != 200:
                raise HTTPException(status_code=response.status_code, detail=f"Embedding API Error: {response.text}")
            
            data = response.json()
            embeddings = data.get("embeddings", [])
            if len(embeddings) < 2:
                raise HTTPException(status_code=500, detail="Failed to retrieve complete vectors")
                
            v1 = embeddings[0]["values"]
            v2 = embeddings[1]["values"]
            sim = cosine_similarity(v1, v2)
            normalized_score = max(0, min(100, int((sim - 0.6) * (100 / 0.35))))
            
            try:
                resumes = client.table('resumes').select('id, content').eq('user_id', user_id).limit(1).execute()
                if resumes.data:
                    content = resumes.data[0].get('content', {})
                    content['ats_used'] = True
                    client.table('resumes').update({"content": content}).eq('id', resumes.data[0]['id']).execute()
            except Exception as e:
                print(f"Failed to update ATS lock: {e}")
            
            return {
                "score": normalized_score,
                "similarity_raw": sim
            }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=503, detail=f"AI connection failed: {str(e)}")


@app.post("/api/v1/ats/analyze")
async def analyze_ats(request: AtsAnalyzeRequest, auth_context: dict = Depends(check_usage_limits)):
    client = auth_context["client"]
    user_id = auth_context["user_id"]
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={GEMINI_API_KEY}"
    
    resume_text = str(request.resume_data.model_dump())[:15000]
    jd_text = request.job_description[:15000]
    
    payload = {
        "contents": [{ 
            "parts": [{ "text": f"Analyze this resume against the Job Description.\nJD:\n{jd_text}\nRESUME:\n{resume_text}" }] 
        }],
        "systemInstruction": { "parts": [{ "text": "You are an ATS. Return JSON strictly matching the schema." }] },
        "generationConfig": {
            "responseMimeType": "application/json",
            "responseSchema": {
                "type": "OBJECT",
                "properties": {
                    "score": { "type": "INTEGER" },
                    "matchedSkills": { "type": "ARRAY", "items": { "type": "STRING" } },
                    "missingSkills": { "type": "ARRAY", "items": { "type": "STRING" } },
                    "optimizationTips": { "type": "ARRAY", "items": { "type": "STRING" } }
                },
                "required": ["score", "matchedSkills", "missingSkills", "optimizationTips"]
            }
        }
    }

    try:
        max_retries = 3
        base_delay = 2 

        async with httpx.AsyncClient() as client_http:
            for attempt in range(max_retries):
                response = await client_http.post(url, headers={"Content-Type": "application/json"}, json=payload, timeout=60.0)
                
                if response.status_code == 200:
                    break
                    
                if response.status_code == 429:
                    if attempt < max_retries - 1:
                        await asyncio.sleep(base_delay * (2 ** attempt))
                        continue
                    else:
                        raise HTTPException(status_code=429, detail="AI Service is currently busy. Please wait and try again.")
                
                raise HTTPException(status_code=response.status_code, detail=f"AI Error: {response.text}")
            
            result = response.json()
            try:
                text_resp = result['candidates'][0]['content']['parts'][0]['text']
                parsed_result = json.loads(text_resp)
                
                try:
                    resumes = client.table('resumes').select('id, content').eq('user_id', user_id).limit(1).execute()
                    if resumes.data:
                        content = resumes.data[0].get('content', {})
                        content['ats_used'] = True
                        client.table('resumes').update({"content": content}).eq('id', resumes.data[0]['id']).execute()
                except Exception as e:
                    print(f"Failed to update ATS lock: {e}")
                    
                return parsed_result
            except Exception:
                raise HTTPException(status_code=500, detail="Failed to parse AI response")
    except HTTPException:
        raise
    except httpx.RequestError as exc:
        raise HTTPException(status_code=503, detail=f"Connection to AI service failed: {str(exc)}")